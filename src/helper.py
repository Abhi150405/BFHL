import os
import io
import requests
import tempfile
from urllib.parse import urlparse
from typing import List
import pandas as pd
import zipfile
from pathlib import Path
import concurrent.futures

# LangChain and Google Generative AI components
from langchain.schema import Document
from langchain_community.document_loaders import (
    PyMuPDFLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredPowerPointLoader,
)
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai

# Libraries for file processing
import pptx
from PIL import Image

# --- CONFIGURATION ---
IMAGE_EXTENSIONS = ['.png', '.jpg', '.jpeg', '.webp', '.gif', '.bmp', '.tiff']

# --- AUTHENTICATION AND INITIALIZATION ---

def get_api_key() -> str:
    """Retrieves the Google API key from environment variables."""
    api_key =  os.getenv("GOOGLE_API_KEY1")
    if not api_key:
        raise ValueError("GOOGLE_API_KEY or GOOGLE_API_KEY1 environment variable not set.")
    return api_key

def initialize_gemini_embeddings() -> GoogleGenerativeAIEmbeddings:
    """Initializes and returns the Google Generative AI embeddings model."""
    return GoogleGenerativeAIEmbeddings(model="models/text-embedding-004", google_api_key=get_api_key())

def initialize_vision_model() -> genai.GenerativeModel:
    """Initializes the Gemini Vision model for image analysis."""
    api_key = get_api_key()
    genai.configure(api_key=api_key)
    return genai.GenerativeModel('gemini-1.5-pro-latest')

# --- IMAGE AND FILE HANDLERS ---

def process_image_with_gemini_vision(image_bytes: bytes, source_description: str) -> Document:
    """Analyzes an image using Gemini Vision and returns its description as a Document."""
    print(f"  -> Analyzing image from '{source_description}' with Gemini Vision...")
    vision_model = initialize_vision_model()
    try:
        img = Image.open(io.BytesIO(image_bytes))
        prompt = "Describe this image in detail. If it contains text, extract it verbatim. If it's a chart or graph, explain its meaning and data points. If it is a technical diagram, explain its components and flow."
        response = vision_model.generate_content([prompt, img])
        return Document(
            page_content=f"Image Description from {source_description}:\n{response.text}",
            metadata={'source': source_description, 'content_type': 'image_analysis', 'original_format': img.format}
        )
    except Exception as e:
        print(f"    - Could not process image from '{source_description}': {e}")
        return Document(page_content=f"[Image from {source_description} could not be processed]", metadata={'source': source_description, 'error': str(e)})

def _process_single_slide(args) -> List[Document]:
    """Internal helper to process one slide, extracting text and images."""
    slide, slide_index, source_filename = args
    slide_source_id = f"{source_filename} (Slide {slide_index + 1})"
    slide_docs = []
    try:
        slide_text = "".join(shape.text.strip() + "\n\n" for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text_frame and shape.text.strip())
        if slide_text:
            slide_docs.append(Document(page_content=slide_text, metadata={'source': slide_source_id, 'content_type': 'slide_text'}))
    except Exception as e:
        print(f"    - Error extracting text from {slide_source_id}: {e}")
    try:
        for shape in slide.shapes:
            if isinstance(shape, pptx.shapes.picture.Picture):
                image_bytes = shape.image.blob
                image_doc = process_image_with_gemini_vision(image_bytes, slide_source_id)
                slide_docs.append(image_doc)
    except Exception as e:
        print(f"    - Error extracting images from {slide_source_id}: {e}")
    return slide_docs

def process_pptx_parallel(file_path: str, source: str) -> List[Document]:
    """Processes a PowerPoint file by handling each slide in parallel."""
    print(f"Processing PowerPoint file '{source}' in parallel...")
    all_documents = []
    try:
        presentation = pptx.Presentation(file_path)
        tasks = [(slide, i, source) for i, slide in enumerate(presentation.slides)]
        if not tasks:
            return []
        print(f"Starting parallel processing for {len(tasks)} slides...")
        with concurrent.futures.ThreadPoolExecutor() as executor:
            results = executor.map(_process_single_slide, tasks)
            for slide_docs in results:
                all_documents.extend(slide_docs)
    except Exception as e:
        print(f"Failed to process PowerPoint file '{source}': {e}")
    return all_documents

def process_csv_as_documents(file_path: str, source: str) -> List[Document]:
    """Processes a CSV file into a single Document with Markdown formatting."""
    print(f"Processing CSV file '{source}' as a single Markdown document...")
    try:
        df = pd.read_csv(file_path)
        if df.empty: return []
        markdown_content = df.to_markdown(index=False)
        return [Document(page_content=markdown_content, metadata={'source': source, 'content_type': 'tabular_data_full'})]
    except Exception as e:
        print(f"Could not process CSV file '{source}': {e}")
        return []

def process_excel_as_documents(file_path: str, source: str) -> List[Document]:
    """Processes an Excel file, creating one Document per sheet."""
    print(f"Processing Excel file '{source}', creating one document per sheet...")
    documents = []
    try:
        with pd.ExcelFile(file_path) as xls:
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                if not df.empty:
                    markdown_content = df.to_markdown(index=False)
                    documents.append(Document(page_content=markdown_content, metadata={'source': source, 'sheet_name': sheet_name, 'content_type': 'tabular_data_full'}))
        return documents
    except Exception as e:
        print(f"Could not process Excel file '{source}': {e}")
        return []

def process_zip_as_documents(zip_path: str, source: str) -> List[Document]:
    """Extracts a ZIP archive and recursively processes each file within."""
    print(f"Processing ZIP archive '{source}'...")
    all_documents = []
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
            for root, _, files in os.walk(temp_dir):
                for file_name in files:
                    local_file_path = os.path.join(root, file_name)
                    print(f"  -> Processing extracted file: {file_name}")
                    docs_from_file = _process_local_file(local_file_path, source=f"{source}/{file_name}")
                    all_documents.extend(docs_from_file)
        except Exception as e:
            print(f"An error occurred while processing ZIP file '{source}': {e}")
    return all_documents

def process_bin_as_documents(file_path: str, source: str) -> List[Document]:
    """Handles a generic .bin file by attempting to read it as an image, then as text."""
    print(f"Handling generic binary file '{source}'...")
    try:
        with Image.open(file_path) as img:
            content = f"Binary file '{source}' identified as image; format '{img.format}', dimensions {img.size}."
            return [Document(page_content=content, metadata={'source': source, 'content_type': 'binary_image'})]
    except Exception:
        try:
            with open(file_path, 'r', errors='ignore') as f:
                content = f.read()
            if content.strip():
                return [Document(page_content=content, metadata={'source': source, 'format': 'binary_text'})]
        except Exception as e:
            print(f"Could not read binary file '{source}' as text: {e}")
    return [Document(page_content=f"[Unreadable binary content in file '{source}']", metadata={'source': source})]

# --- Internal Core Processing Logic ---

def _process_local_file(file_path: str, source: str) -> List[Document]:
    """Internal function to process a file from a local path based on its extension."""
    file_ext = Path(file_path).suffix.lower()
    print(f"Dispatching local file '{source}' with extension '{file_ext}' to handler.")
    if file_ext == '.pdf':
        return PyMuPDFLoader(file_path=file_path).load()
    elif file_ext == '.docx':
        return UnstructuredWordDocumentLoader(file_path=file_path).load()
    elif file_ext == '.csv':
        return process_csv_as_documents(file_path, source=source)
    elif file_ext in ['.xls', '.xlsx']:
        return process_excel_as_documents(file_path, source=source)
    elif file_ext == '.pptx':
        return process_pptx_parallel(file_path, source=source)
    elif file_ext in IMAGE_EXTENSIONS:
        with open(file_path, 'rb') as f:
            return [process_image_with_gemini_vision(f.read(), source)]
    elif file_ext == '.zip':
        return process_zip_as_documents(file_path, source=source)
    elif file_ext == '.bin':
        return process_bin_as_documents(file_path, source=source)
    else:
        print(f"Warning: No specific handler for file extension '{file_ext}'. No documents loaded.")
        return []

# --- Main Document Loading Orchestrator ---

def load_doc_from_url(url: str) -> List[Document]:
    """Main public function. Downloads a document from a URL and processes it."""
    file_name = os.path.basename(urlparse(url).path)
    file_ext = os.path.splitext(file_name)[1].lower()

    if not file_ext:
        print(f"URL has no file extension. Attempting to infer from Content-Type header...")
        try:
            head_response = requests.head(url, allow_redirects=True, timeout=20)
            head_response.raise_for_status()
            content_type = head_response.headers.get('Content-Type', '').lower()
            ext_map = {'.csv': 'csv', '.pdf': 'pdf', '.zip': 'zip', '.xlsx': 'spreadsheetml', '.docx': 'wordprocessingml', '.pptx': 'presentationml'}
            for ext, type_str in ext_map.items():
                if type_str in content_type:
                    file_ext = ext
                    break
            if file_ext:
                print(f"Inferred file extension '{file_ext}' from Content-Type: {content_type}")
                if '.' not in file_name: file_name += file_ext
            else:
                print(f"Could not infer a specific extension from Content-Type: {content_type}")
        except requests.exceptions.RequestException as e:
            print(f"Warning: Could not perform HEAD request: {e}")

    temp_file_path = None
    try:
        print(f"Downloading document from {url}...")
        response = requests.get(url, timeout=60)
        response.raise_for_status()
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext or ".tmp") as temp_file:
            temp_file.write(response.content)
            temp_file_path = temp_file.name
        
        pages = _process_local_file(temp_file_path, source=file_name or url)
        print(f"Document processing complete. Loaded {len(pages)} total documents/pages.")
        return pages
    except requests.exceptions.RequestException as e:
        print(f"Failed to download document from URL: {e}")
        return []
    except Exception as e:
        print(f"An error occurred during document processing: {e}")
        return []
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.unlink(temp_file_path)
            print(f"Cleaned up temporary file: {temp_file_path}")
